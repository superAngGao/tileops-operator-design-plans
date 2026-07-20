#!/usr/bin/env python3
from __future__ import annotations

import csv
import json
import re
import subprocess
from collections import Counter, defaultdict
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path

import matplotlib as mpl
import matplotlib.pyplot as plt
import pandas as pd
from matplotlib import font_manager
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


REPO = Path("/home/ga/TileOPs")
OUT = Path("/home/ga/tileops-gqa-plan/superAngGao-tileops-contrib")
PR_JSON = OUT / "superAngGao_prs.json"
AUTHOR_EMAIL = "gaoang0125@163.com"
REMOTE_REF = "upstream/main"


FONT_PATH = Path("/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf")
if FONT_PATH.exists():
    font_manager.fontManager.addfont(str(FONT_PATH))
    CJK_FONT = font_manager.FontProperties(fname=str(FONT_PATH)).get_name()
else:
    CJK_FONT = "DejaVu Sans"

mpl.rcParams["font.family"] = ["DejaVu Sans", CJK_FONT]
mpl.rcParams["svg.fonttype"] = "path"
mpl.rcParams["axes.unicode_minus"] = False
plt.rcParams["figure.facecolor"] = "white"


COLORS = {
    "ink": "#172026",
    "muted": "#63707a",
    "blue": "#2f6fed",
    "teal": "#00a896",
    "amber": "#f5a524",
    "red": "#e25555",
    "violet": "#7c5cff",
    "green": "#2fb344",
    "line": "#dce4ea",
    "panel": "#f7f9fb",
}


def run(cmd: list[str]) -> str:
    return subprocess.check_output(cmd, cwd=REPO, text=True).strip()


def fetch_inputs() -> None:
    run(["git", "fetch", "upstream", "main", "--prune"])
    pr_json = subprocess.check_output(
        [
            "gh",
            "pr",
            "list",
            "-R",
            "tile-ai/TileOPs",
            "--author",
            "superAngGao",
            "--state",
            "all",
            "--limit",
            "300",
            "--json",
            "number,title,state,isDraft,createdAt,updatedAt,closedAt,mergedAt,"
            "additions,deletions,changedFiles,url,baseRefName,headRefName",
        ],
        text=True,
    )
    PR_JSON.write_text(pr_json, encoding="utf-8")


@dataclass
class Commit:
    sha: str
    date: str
    author: str
    email: str
    subject: str


def load_commits() -> list[Commit]:
    raw = run(
        [
            "git",
            "log",
            REMOTE_REF,
            f"--author={AUTHOR_EMAIL}",
            "--format=%H%x09%ad%x09%an%x09%ae%x09%s",
            "--date=short",
        ]
    )
    commits = []
    for line in raw.splitlines():
        sha, date, author, email, subject = line.split("\t", 4)
        commits.append(Commit(sha, date, author, email, subject))
    return commits


def load_numstat() -> tuple[Counter[str], Counter[str], Counter[str], int, int, int]:
    raw = run(
        [
            "git",
            "log",
            REMOTE_REF,
            f"--author={AUTHOR_EMAIL}",
            "--numstat",
            "--format=COMMIT%x09%H%x09%ad%x09%s",
            "--date=short",
        ]
    )
    module_lines: Counter[str] = Counter()
    file_lines: Counter[str] = Counter()
    area_files: Counter[str] = Counter()
    additions = deletions = changed_entries = 0
    for line in raw.splitlines():
        if not line or line.startswith("COMMIT\t"):
            continue
        parts = line.split("\t")
        if len(parts) < 3:
            continue
        add_s, del_s, path = parts[0], parts[1], parts[2]
        if add_s == "-" or del_s == "-":
            continue
        add, delete = int(add_s), int(del_s)
        total = add + delete
        additions += add
        deletions += delete
        changed_entries += 1
        top = path.split("/", 1)[0]
        module_lines[top] += total
        file_lines[path] += total
        area_files[top] += 1
    return module_lines, file_lines, area_files, additions, deletions, changed_entries


def classify_subject(subject: str) -> str:
    m = re.match(r"^\[([^\]]+)\]", subject)
    if not m:
        return "Other"
    first = m.group(1).split("]")[0].split("[")[0].split("/", 1)[0].split("][", 1)[0]
    return first


def save_csv(name: str, rows: list[dict]) -> None:
    if not rows:
        return
    with (OUT / name).open("w", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=list(rows[0].keys()))
        writer.writeheader()
        writer.writerows(rows)


def style_axes(ax) -> None:
    ax.spines[["top", "right", "left"]].set_visible(False)
    ax.spines["bottom"].set_color(COLORS["line"])
    ax.tick_params(colors=COLORS["muted"])
    ax.grid(axis="y", color=COLORS["line"], linewidth=0.8, alpha=0.8)
    ax.set_axisbelow(True)


def savefig(fig, stem: str) -> None:
    fig.savefig(OUT / f"{stem}.png", dpi=220, bbox_inches="tight")
    fig.savefig(OUT / f"{stem}.svg", bbox_inches="tight")
    plt.close(fig)


def make_dashboard(summary: dict, monthly: pd.DataFrame, top_authors: pd.DataFrame) -> None:
    fig = plt.figure(figsize=(16, 9))
    gs = fig.add_gridspec(3, 4, height_ratios=[0.9, 1.25, 1.25], hspace=0.62, wspace=0.42)
    fig.suptitle("superAngGao / Ang Gao 在 TileOPs 的贡献概览", x=0.06, y=0.97, ha="left",
                 fontsize=27, fontweight="bold", color=COLORS["ink"])
    fig.text(0.06, 0.915, f"统计口径：tile-ai/TileOPs {REMOTE_REF}，更新至 {summary['latest_commit_date']}",
             fontsize=12.5, color=COLORS["muted"])

    kpis = [
        ("主线提交", f"{summary['main_commits']}", "merged commits"),
        ("作者排名", f"#{summary['author_rank']}", f"{summary['contribution_share']:.1f}% of commits"),
        ("Merged PR", f"{summary['merged_prs']}", f"{summary['all_prs']} total PRs"),
        ("代码变更", f"{summary['git_additions'] + summary['git_deletions']:,}", "additions + deletions"),
    ]
    for i, (label, value, sub) in enumerate(kpis):
        ax = fig.add_subplot(gs[0, i])
        ax.set_facecolor(COLORS["panel"])
        for sp in ax.spines.values():
            sp.set_visible(False)
        ax.set_xticks([])
        ax.set_yticks([])
        ax.text(0.06, 0.72, label, transform=ax.transAxes, fontsize=14, color=COLORS["muted"])
        ax.text(0.06, 0.28, value, transform=ax.transAxes, fontsize=31, fontweight="bold", color=COLORS["ink"])
        ax.text(0.06, 0.08, sub, transform=ax.transAxes, fontsize=11, color=COLORS["muted"])

    ax1 = fig.add_subplot(gs[1:, :2])
    monthly = monthly.copy()
    monthly["cumulative_commits"] = monthly["commits"].cumsum()
    ax1.fill_between(monthly["month"], monthly["cumulative_commits"], color=COLORS["blue"], alpha=0.16)
    ax1.plot(monthly["month"], monthly["cumulative_commits"], color=COLORS["blue"], linewidth=3, marker="o", markersize=8)
    for x, y in zip(monthly["month"], monthly["cumulative_commits"]):
        ax1.annotate(str(y), (x, y), textcoords="offset points", xytext=(0, 8),
                     ha="center", fontsize=11, color=COLORS["ink"])
    ax1.set_title("主线提交累计趋势", loc="left", fontsize=16, fontweight="bold", color=COLORS["ink"])
    ax1.set_ylabel("Commits", color=COLORS["muted"])
    style_axes(ax1)

    ax2 = fig.add_subplot(gs[1:, 2:])
    own = summary["main_commits"]
    others = summary["total_main_commits"] - own
    ax2.pie(
        [own, others],
        labels=["superAngGao", "其他贡献"],
        startangle=90,
        counterclock=False,
        colors=[COLORS["teal"], "#e6edf3"],
        autopct=lambda p: f"{p:.1f}%",
        pctdistance=0.74,
        wedgeprops={"linewidth": 2, "edgecolor": "white"},
        textprops={"fontsize": 13, "color": COLORS["ink"]},
    )
    ax2.text(0, 0.07, f"{own}", ha="center", va="center", fontsize=36,
             fontweight="bold", color=COLORS["ink"])
    ax2.text(0, -0.13, f"/ {summary['total_main_commits']} commits", ha="center",
             va="center", fontsize=12, color=COLORS["muted"])
    ax2.set_title("主线提交贡献占比", loc="left", fontsize=16, fontweight="bold", color=COLORS["ink"])
    savefig(fig, "00_dashboard")


def make_charts(summary: dict, monthly: pd.DataFrame, categories: pd.DataFrame,
                modules: pd.DataFrame, prs: list[dict]) -> None:
    fig, ax = plt.subplots(figsize=(10.5, 5.6))
    bars = ax.bar(monthly["month"], monthly["commits"], color=COLORS["blue"], width=0.58)
    ax.bar_label(bars, padding=4, fontsize=12)
    ax.set_title("superAngGao 主线提交月度趋势", loc="left", fontsize=18, fontweight="bold")
    ax.set_ylabel("Commits")
    style_axes(ax)
    savefig(fig, "01_monthly_commits")

    fig, ax = plt.subplots(figsize=(8, 8))
    wedges, texts, autotexts = ax.pie(
        categories["commits"],
        labels=categories["category"],
        startangle=90,
        autopct=lambda p: f"{p:.0f}%" if p >= 5 else "",
        colors=[COLORS["blue"], COLORS["teal"], COLORS["amber"], COLORS["violet"], COLORS["green"], "#8aa2b2",
                "#c64f8a", "#4c9f70", "#5c6bc0", "#9e9e9e"][: len(categories)],
        wedgeprops={"linewidth": 1, "edgecolor": "white"},
        textprops={"fontsize": 11, "color": COLORS["ink"]},
    )
    ax.text(0, 0.05, str(summary["main_commits"]), ha="center", va="center",
            fontsize=34, fontweight="bold", color=COLORS["ink"])
    ax.text(0, -0.16, "main commits", ha="center", va="center", fontsize=11, color=COLORS["muted"])
    ax.set_title("提交类型分布", loc="left", fontsize=18, fontweight="bold")
    savefig(fig, "02_commit_type_mix")

    fig, ax = plt.subplots(figsize=(11, 6))
    top = modules.head(10).copy()
    ax.barh(top["module"][::-1], top["changed_lines"][::-1], color=COLORS["teal"])
    ax.set_title("代码变更覆盖模块 Top 10", loc="left", fontsize=18, fontweight="bold")
    ax.set_xlabel("Changed lines (additions + deletions)")
    style_axes(ax)
    ax.grid(axis="x", color=COLORS["line"], linewidth=0.8, alpha=0.8)
    ax.grid(axis="y", visible=False)
    savefig(fig, "03_module_impact")

    status_counter = Counter(pr["state"] for pr in prs)
    pr_status = pd.DataFrame(
        [
            {"status": "Merged", "count": status_counter.get("MERGED", 0), "color": COLORS["green"]},
            {"status": "Open", "count": status_counter.get("OPEN", 0), "color": COLORS["blue"]},
            {"status": "Closed", "count": status_counter.get("CLOSED", 0), "color": COLORS["muted"]},
        ]
    )
    fig, ax = plt.subplots(figsize=(8.2, 6))
    bars = ax.bar(pr_status["status"], pr_status["count"], color=pr_status["color"])
    ax.bar_label(bars, padding=4, fontsize=13)
    ax.set_title("GitHub PR 状态", loc="left", fontsize=18, fontweight="bold")
    ax.set_ylabel("PRs")
    style_axes(ax)
    savefig(fig, "04_pr_status")


def add_title(slide, text: str, subtitle: str | None = None) -> None:
    title = slide.shapes.add_textbox(Inches(0.45), Inches(0.24), Inches(12.4), Inches(0.5))
    p = title.text_frame.paragraphs[0]
    p.text = text
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(25)
    p.font.bold = True
    p.font.color.rgb = RGBColor(23, 32, 38)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.47), Inches(0.76), Inches(12.1), Inches(0.34))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.name = "Microsoft YaHei"
        sp.font.size = Pt(11.5)
        sp.font.color.rgb = RGBColor(99, 112, 122)


def add_metric(slide, x: float, y: float, label: str, value: str, note: str) -> None:
    box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.75), Inches(1.1))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(247, 249, 251)
    box.line.color.rgb = RGBColor(232, 238, 243)
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = label
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(99, 112, 122)
    value_p = frame.add_paragraph()
    value_p.text = value
    value_p.font.name = "Arial"
    value_p.font.size = Pt(26)
    value_p.font.bold = True
    value_p.font.color.rgb = RGBColor(23, 32, 38)
    note_p = frame.add_paragraph()
    note_p.text = note
    note_p.font.name = "Arial"
    note_p.font.size = Pt(9.5)
    note_p.font.color.rgb = RGBColor(99, 112, 122)


def make_pptx(summary: dict) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    add_title(
        slide,
        "superAngGao 在 TileOPs 的贡献统计",
        f"tile-ai/TileOPs · {REMOTE_REF} · 更新至 {summary['latest_commit_date']}",
    )
    add_metric(slide, 0.55, 1.25, "主线提交", str(summary["main_commits"]), "merged commits")
    add_metric(slide, 3.75, 1.25, "作者排名", f"#{summary['author_rank']}", f"{summary['contribution_share']:.1f}% of commits")
    add_metric(slide, 6.95, 1.25, "Merged PR", str(summary["merged_prs"]), f"{summary['all_prs']} total PRs")
    add_metric(
        slide,
        10.15,
        1.25,
        "代码变更",
        f"{summary['git_additions'] + summary['git_deletions']:,}",
        "additions + deletions",
    )
    slide.shapes.add_picture(str(OUT / "01_monthly_commits.png"), Inches(0.55), Inches(2.75), width=Inches(5.95))
    slide.shapes.add_picture(str(OUT / "04_pr_status.png"), Inches(7.0), Inches(2.78), width=Inches(4.85))

    slide = prs.slides.add_slide(blank)
    add_title(slide, "提交趋势与类型分布", "突出 2026-03 的集中合入，以及 Fix/CI/Feat 的主要贡献方向")
    slide.shapes.add_picture(str(OUT / "01_monthly_commits.png"), Inches(0.55), Inches(1.25), width=Inches(6.15))
    slide.shapes.add_picture(str(OUT / "02_commit_type_mix.png"), Inches(7.25), Inches(1.05), width=Inches(5.0))

    slide = prs.slides.add_slide(blank)
    add_title(slide, "模块影响范围与 PR 流程", "主线代码、benchmark、tests 和 CI 都有覆盖，可展示为工程闭环贡献")
    slide.shapes.add_picture(str(OUT / "03_module_impact.png"), Inches(0.55), Inches(1.25), width=Inches(6.45))
    slide.shapes.add_picture(str(OUT / "04_pr_status.png"), Inches(7.25), Inches(1.45), width=Inches(4.85))

    slide = prs.slides.add_slide(blank)
    add_title(slide, "一页总览图", "可作为汇报中的独立截图页，也可拆分为 KPI + 趋势 + 排名")
    slide.shapes.add_picture(str(OUT / "00_dashboard.png"), Inches(0.3), Inches(1.03), width=Inches(12.72))

    gqa_perf_chart = OUT / "05_gqa_perf_pr871.png"
    if gqa_perf_chart.exists():
        slide = prs.slides.add_slide(blank)
        add_title(slide, "GQA Dense Prefill 性能对比", "数据来自 nightly run 25259450950：causal inference prefill，相对 FA3 归一化")
        slide.shapes.add_picture(str(gqa_perf_chart), Inches(0.25), Inches(0.95), width=Inches(12.85))

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.LEFT

    prs.save(OUT / "superAngGao_tileops_contribution.pptx")


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    fetch_inputs()

    commits = load_commits()
    prs = json.loads(PR_JSON.read_text(encoding="utf-8"))
    module_lines, file_lines, area_files, additions, deletions, changed_entries = load_numstat()

    total_commits = int(run(["git", "rev-list", "--count", REMOTE_REF]))
    latest_commit_date = run(["git", "log", "-1", REMOTE_REF, "--format=%ad", "--date=short"])
    first_commit_date = commits[-1].date if commits else ""
    author_counts = []
    for line in run(["git", "shortlog", "-sne", REMOTE_REF]).splitlines():
        count_s, author = line.strip().split("\t", 1)
        author_counts.append({"commits": int(count_s), "author": author})
    top_authors = pd.DataFrame(author_counts)
    top_authors["author_short"] = top_authors["author"].str.replace(r"\s+<.*", "", regex=True)
    author_rank = int(top_authors.index[top_authors["author"].str.contains(AUTHOR_EMAIL)][0]) + 1

    monthly_counter = Counter(c.date[:7] for c in commits)
    monthly = pd.DataFrame(
        [{"month": month, "commits": monthly_counter[month]} for month in sorted(monthly_counter)]
    )
    categories_counter = Counter(classify_subject(c.subject) for c in commits)
    categories = pd.DataFrame(
        [{"category": k, "commits": v} for k, v in categories_counter.most_common()]
    )
    modules = pd.DataFrame(
        [{"module": k, "changed_lines": v, "file_entries": area_files[k]} for k, v in module_lines.most_common()]
    )

    summary = {
        "repo": "tile-ai/TileOPs",
        "ref": REMOTE_REF,
        "latest_commit_date": latest_commit_date,
        "first_main_commit_date": first_commit_date,
        "main_commits": len(commits),
        "total_main_commits": total_commits,
        "contribution_share": len(commits) / total_commits * 100,
        "author_rank": author_rank,
        "all_prs": len(prs),
        "merged_prs": sum(1 for pr in prs if pr["mergedAt"]),
        "open_prs": sum(1 for pr in prs if pr["state"] == "OPEN"),
        "closed_unmerged_prs": sum(1 for pr in prs if pr["state"] == "CLOSED"),
        "draft_prs": sum(1 for pr in prs if pr["isDraft"]),
        "pr_additions": sum(pr["additions"] for pr in prs),
        "pr_deletions": sum(pr["deletions"] for pr in prs),
        "pr_changed_files": sum(pr["changedFiles"] for pr in prs),
        "git_additions": additions,
        "git_deletions": deletions,
        "git_changed_file_entries": changed_entries,
        "top_commit_month": monthly.sort_values("commits", ascending=False).iloc[0].to_dict(),
        "top_modules": modules.head(5).to_dict("records"),
        "top_categories": categories.head(5).to_dict("records"),
    }

    (OUT / "summary.json").write_text(json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8")
    save_csv("main_commits.csv", [c.__dict__ for c in commits])
    save_csv("monthly_commits.csv", monthly.to_dict("records"))
    save_csv("commit_categories.csv", categories.to_dict("records"))
    save_csv("module_impact.csv", modules.to_dict("records"))
    save_csv("top_files.csv", [{"file": k, "changed_lines": v} for k, v in file_lines.most_common(25)])
    save_csv("top_authors.csv", top_authors.to_dict("records"))
    save_csv(
        "prs.csv",
        [
            {
                "number": pr["number"],
                "state": pr["state"],
                "mergedAt": pr["mergedAt"],
                "createdAt": pr["createdAt"],
                "title": pr["title"],
                "additions": pr["additions"],
                "deletions": pr["deletions"],
                "changedFiles": pr["changedFiles"],
                "url": pr["url"],
            }
            for pr in prs
        ],
    )

    make_dashboard(summary, monthly, top_authors)
    make_charts(summary, monthly, categories, modules, prs)
    make_pptx(summary)

    report = f"""# superAngGao 在 TileOPs 的贡献统计

统计时间：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}，仓库：tile-ai/TileOPs，口径：`{REMOTE_REF}`。

## 可直接放 PPT 的结论

- 主线贡献：{summary['main_commits']} 个 commits，占主线 {summary['total_main_commits']} 个 commits 的 {summary['contribution_share']:.1f}%，作者排名第 {summary['author_rank']}。
- PR 贡献：GitHub 上由 `superAngGao` 发起 {summary['all_prs']} 个 PR，其中 {summary['merged_prs']} 个已合入，{summary['open_prs']} 个仍打开。
- 代码规模：按主线 commit numstat 统计，新增 {summary['git_additions']:,} 行、删除 {summary['git_deletions']:,} 行，合计 {summary['git_additions'] + summary['git_deletions']:,} 行变更。
- 高峰月份：{summary['top_commit_month']['month']}，主线合入 {int(summary['top_commit_month']['commits'])} 个 commits。
- 主要类型：{', '.join(f"{r['category']} {r['commits']}" for r in summary['top_categories'])}。
- 覆盖模块：{', '.join(f"{r['module']} {r['changed_lines']:,} 行" for r in summary['top_modules'])}。

## PPT 图表文件

- `00_dashboard.png/svg`：一页总览仪表盘。
- `01_monthly_commits.png/svg`：月度主线提交趋势。
- `02_commit_type_mix.png/svg`：提交类型分布。
- `03_module_impact.png/svg`：模块影响范围。
- `04_pr_status.png/svg`：PR 状态统计。
- `superAngGao_tileops_contribution.pptx`：已排版好的 4 页 16:9 PPT 草稿。

## 口径说明

- Git 主线统计按 author email `{AUTHOR_EMAIL}` 归并，即 `Ang Gao <{AUTHOR_EMAIL}>`。
- PR 统计来自 GitHub CLI：`gh pr list -R tile-ai/TileOPs --author superAngGao --state all`。
- 行数统计来自 `git log {REMOTE_REF} --author={AUTHOR_EMAIL} --numstat`，适合展示贡献规模，不等同于当前仓库净增行数。
"""
    (OUT / "README.md").write_text(report, encoding="utf-8")
    print(json.dumps(summary, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
